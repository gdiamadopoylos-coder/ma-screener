# ppt_export.py
from __future__ import annotations

import io
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple, Union

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE_FILENAME = "merjure_template_v1.pptx"
TEMPLATE_PATH = Path(__file__).resolve().parent / TEMPLATE_FILENAME


def _load_template() -> Presentation:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(
            f"Template not found at: {TEMPLATE_PATH}\n"
            f"Put '{TEMPLATE_FILENAME}' in the same folder as ppt_export.py."
        )
    return Presentation(str(TEMPLATE_PATH))


def _layout(prs: Presentation, preferred: Sequence[int]) -> Any:
    """Return the first available slide layout from 'preferred', else the last, else 0."""
    n = len(prs.slide_layouts)
    if n == 0:
        raise RuntimeError("This PPTX has zero slide layouts (corrupt template).")
    for idx in preferred:
        if 0 <= idx < n:
            return prs.slide_layouts[idx]
    # safe fallbacks
    return prs.slide_layouts[min(n - 1, 0)]


def _safe_str(x: Any, default: str = "") -> str:
    if x is None:
        return default
    s = str(x).strip()
    return s if s else default


def _safe_list_str(x: Any) -> List[str]:
    if x is None:
        return []
    if isinstance(x, (list, tuple)):
        return [str(v).strip() for v in x if v is not None and str(v).strip()]
    s = str(x).strip()
    return [s] if s else []


def _today_str() -> str:
    return date.today().isoformat()


def _find_title_placeholder(slide) -> Optional[Any]:
    for shp in slide.placeholders:
        if getattr(shp, "has_text_frame", False) and getattr(shp.placeholder_format, "idx", None) == 0:
            return shp
    for shp in slide.placeholders:
        if getattr(shp, "has_text_frame", False):
            return shp
    return None


def _find_body_placeholder(slide) -> Optional[Any]:
    for shp in slide.placeholders:
        if getattr(shp, "has_text_frame", False) and getattr(shp.placeholder_format, "idx", None) == 1:
            return shp
    for shp in slide.placeholders:
        if getattr(shp, "has_text_frame", False) and getattr(shp.placeholder_format, "idx", None) != 0:
            return shp
    return None


def _set_text(shape, text: str, font_size: int = 24, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold


def _set_bullets(shape, bullets: Sequence[str], font_size: int = 18) -> None:
    tf = shape.text_frame
    tf.clear()
    clean = [b.strip() for b in bullets if b and str(b).strip()]
    if not clean:
        p = tf.paragraphs[0]
        p.text = ""
        p.font.size = Pt(font_size)
        return

    p0 = tf.paragraphs[0]
    p0.text = clean[0]
    p0.level = 0
    p0.font.size = Pt(font_size)

    for b in clean[1:]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)


def _add_textbox(slide, left, top, width, height, text: str, font_size: int = 18, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold


def _add_bullets_box(slide, left, top, width, height, title: str, bullets: Sequence[str]) -> None:
    if title:
        _add_textbox(slide, left, top, width, Inches(0.4), title, font_size=20, bold=True)
        top2 = top + Inches(0.45)
        h2 = height - Inches(0.45)
    else:
        top2 = top
        h2 = height

    box = slide.shapes.add_textbox(left, top2, width, h2)
    _set_bullets(box, bullets, font_size=16)


def _add_table(slide, title: str, columns: Sequence[str], rows: Sequence[Sequence[Any]]) -> None:
    _add_textbox(slide, Inches(0.7), Inches(0.6), Inches(12.0), Inches(0.5), title, font_size=22, bold=True)

    n_rows = 1 + len(rows)
    n_cols = max(1, len(columns))

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.9), Inches(1.6),
        Inches(11.6), Inches(4.8)
    )
    table = table_shape.table

    for j, col in enumerate(columns):
        if j >= n_cols:
            break
        cell = table.cell(0, j)
        cell.text = _safe_str(col)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows, start=1):
        for j in range(n_cols):
            val = ""
            if row is not None and j < len(row):
                val = _safe_str(row[j])
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER


def _add_cover_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(_layout(prs, [0, 1, 2]))
    title_ph = _find_title_placeholder(slide)
    body_ph = _find_body_placeholder(slide)

    if title_ph:
        _set_text(title_ph, title, font_size=40, bold=True, align=PP_ALIGN.CENTER)
    else:
        _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(12.0), Inches(0.8), title, font_size=40, bold=True, align=PP_ALIGN.CENTER)

    if body_ph:
        _set_text(body_ph, subtitle, font_size=20, bold=False, align=PP_ALIGN.CENTER)
    else:
        _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(12.0), Inches(0.6), subtitle, font_size=20, bold=False, align=PP_ALIGN.CENTER)


def _add_exec_summary(prs: Presentation, thesis: List[str], synergies: List[str], risks: List[str]) -> None:
    slide = prs.slides.add_slide(_layout(prs, [1, 0, 2, 3]))
    title_ph = _find_title_placeholder(slide)

    if title_ph:
        _set_text(title_ph, "Executive Summary", font_size=28, bold=True)
    else:
        _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.5), "Executive Summary", font_size=28, bold=True)

    col_w = Inches(4.0)
    top = Inches(1.3)
    height = Inches(5.5)

    _add_bullets_box(slide, Inches(0.7), top, col_w, height, "Merger Thesis", thesis)
    _add_bullets_box(slide, Inches(4.85), top, col_w, height, "Synergies (Hypotheses)", synergies)
    _add_bullets_box(slide, Inches(9.0), top, col_w, height, "Key Risks", risks)


def _add_kpi_table(prs: Presentation, kpis: Any) -> None:
    slide = prs.slides.add_slide(_layout(prs, [5, 1, 0, 2]))
    rows: List[List[str]] = []

    if isinstance(kpis, dict):
        rows = [[_safe_str(k), _safe_str(v)] for k, v in kpis.items()]
    elif isinstance(kpis, list):
        for item in kpis:
            if isinstance(item, (list, tuple)) and len(item) >= 2:
                rows.append([_safe_str(item[0]), _safe_str(item[1])])

    if not rows:
        rows = [["EV / EBITDA", ""], ["Revenue", ""], ["EBITDA margin", ""], ["Net debt / EBITDA", ""]]

    _add_table(slide, "Key Metrics Snapshot", ["Metric", "Value"], rows)


def _add_sources(prs: Presentation, sources: List[str]) -> None:
    slide = prs.slides.add_slide(_layout(prs, [1, 0, 2, 3]))
    title_ph = _find_title_placeholder(slide)
    body_ph = _find_body_placeholder(slide)

    if title_ph:
        _set_text(title_ph, "Sources / Notes", font_size=28, bold=True)
    else:
        _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.5), "Sources / Notes", font_size=28, bold=True)

    if body_ph:
        _set_bullets(body_ph, sources, font_size=18)
    else:
        _add_bullets_box(slide, Inches(0.9), Inches(1.4), Inches(11.2), Inches(5.4), "", sources)


def build_merger_pptx(payload: Dict[str, Any]) -> bytes:
    prs = _load_template()

    # enforce 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    acquirer = _safe_str(payload.get("acquirer"), "Acquirer")
    target = _safe_str(payload.get("target"), "Target")
    date_str = _safe_str(payload.get("date"), _today_str())

    subtitle = f"{acquirer} + Top Ranked Targets  |  {date_str}"
    cover_title = "M&A Intelligence — Merger Candidate"

    thesis = _safe_list_str(payload.get("thesis_bullets")) or [
        "Strategic adjacency with clear value-creation path.",
        "Capability expansion and portfolio optimisation.",
        "Defensible rationale under current scenario assumptions.",
    ]
    synergies = _safe_list_str(payload.get("synergies_bullets")) or [
        "SG&A rationalisation (duplicate functions).",
        "Procurement scale benefits.",
        "Commercial cross-sell / revenue uplift.",
    ]
    risks = _safe_list_str(payload.get("risks_bullets")) or [
        "Integration complexity and execution risk.",
        "Regulatory / customer concentration sensitivity.",
        "Leverage / financing and macro uncertainty.",
    ]
    sources = _safe_list_str(payload.get("sources_bullets")) or [
        "Input universe file (user supplied CSV/XLSX)",
        "Merjure scoring model (Pilot V1)",
    ]

    _add_cover_slide(prs, cover_title, subtitle)
    _add_exec_summary(prs, thesis=thesis, synergies=synergies, risks=risks)
    _add_kpi_table(prs, kpis=payload.get("kpis"))
    _add_sources(prs, sources=sources)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()